import threading
from queue import Queue
import requests
from bs4 import BeautifulSoup
import time
from urllib.parse import urljoin, urlparse
import hashlib
import psycopg2
import io
import PyPDF2
import docx
import pptx
from datetime import datetime
from selenium import webdriver
import base64
from selenium.webdriver.firefox.options import Options

# Define a lock for thread-safe access to shared resources
lock = threading.Lock()


# Function to retrieve and parse HTML content using requests
def fetch_html(url):
    try:
        firefox_options = Options()
        firefox_options.add_argument("-headless")
        driver = webdriver.Firefox(options=firefox_options)
        driver.get(url)
        response = requests.head(url, timeout=5)
        status_code = response.status_code
        html = driver.page_source
        driver.quit()
        return html, status_code
    except Exception as e:
        print(f"Error fetching {url}: {e}")
    return None


# Function to extract links from HTML content
def extract_links(html, base_url):
    soup = BeautifulSoup(html, 'html.parser')
    links = set()
    for link in soup.find_all('a', href=True):
        href = link.get('href')
        if href.startswith('http') or href.startswith('https'):
            links.add(canonicalize_url(href))
        else:
            links.add(canonicalize_url(urljoin(base_url, href)))
    return links


# Function to respect crawl delay
def respect_crawl_delay(domain, delay=5):
    current_time = time.time()
    last_access_time = last_access_times.get(domain, 0)
    if current_time - last_access_time < delay:
        time.sleep(delay - (current_time - last_access_time))
    last_access_times[domain] = time.time()


# Function to canonicalize URLs
def canonicalize_url(url):
    parsed_url = urlparse(url)
    # Scheme and netloc are converted to lowercase
    scheme = parsed_url.scheme.lower()
    netloc = parsed_url.netloc.lower()
    # Path is stripped of trailing slashes and converted to lowercase
    path = parsed_url.path.rstrip('/').lower()
    # Query and fragment are ignored
    return f"{scheme}://{netloc}{path}"


# Function to crawl a single URL
def crawl_url(url, site_id):
    global current_url
    html, status_code = None, None
    try:
        html, status_code = fetch_html(url)
    except TypeError:
        print("Cannot unpack html and status code, probably, didn't receive a response")
    if html is not None and status_code is not None:  # Store canonicalized URL
        canonical_url = url
        # Check for duplicate content
        if not is_duplicate_html(html):
            if not is_duplicate_url(url):
                # Extract links
                links = extract_links(html, url)
                for link in links:
                    # print("link: " + link)
                    cur.execute("INSERT INTO frontier (link, from_link) VALUES (%s, %s)", (link, current_url))
                    # Store data in the database
                store_data(url, canonical_url, html, status_code, site_id)
            else:
                cur.execute("SELECT id FROM page WHERE url = %s", (from_url,))
                from_id = cur.fetchone()[0]
                cur.execute("SELECT id FROM page WHERE url = %s", (current_url,))
                to_id = cur.fetchone()[0]
                cur.execute("INSERT INTO link (from_page, to_page) VALUES (%s, %s)",
                        (from_id, to_id))
                conn.commit()
    else:
        pass


# Function to check for duplicate content
def is_duplicate_html(html):
    # Check if the hash of the HTML content already exists in the database
    html_hash = hashlib.sha256(html.encode()).hexdigest()
    cur.execute("SELECT * FROM page WHERE html_content_hash = %s", (html_hash,))
    return cur.fetchone() is not None


def is_duplicate_url(url):
    canonical_url = url
    cur.execute("SELECT * FROM page WHERE url = %s", (canonical_url,))
    return cur.fetchone() is not None


# Function to store data in the database
def store_data(url, canonical_url, html, status_code, site_id):
    page_types = ["HTML", "BINARY", "DUPLICATE", "FRONTIER"]
    if '.pdf' in url or '.doc' in url or '.docx' in url or '.ppt' in url or '.pptx' in url:
        html_hash = hashlib.sha256(html.encode()).hexdigest()
        cur.execute("INSERT INTO page (site_id, page_type_code, url, html_content_hash, html_content, http_status_code, accessed_time) "
                    "VALUES (%s, %s, %s, %s, %s, %s, %s) RETURNING id",
                    (site_id, page_types[1], canonical_url, html_hash, html, status_code, datetime.now()))
        page_id = cur.fetchone()[0]  # Get the inserted page ID
        conn.commit()
    else:
        html_hash = hashlib.sha256(html.encode()).hexdigest()
        cur.execute("INSERT INTO page (site_id, page_type_code, url, html_content_hash, html_content, http_status_code, accessed_time) "
                    "VALUES (%s, %s, %s, %s, %s, %s, %s) RETURNING id",
                    (site_id, page_types[0], canonical_url, html_hash, html, status_code, datetime.now()))
        page_id = cur.fetchone()[0]  # Get the inserted page ID
        conn.commit()

    # Check for specific file extensions in the URL and store them in a separate table
    file_extensions = ['.pdf', '.doc', '.docx', '.ppt', '.pptx']
    for ext in file_extensions:
        if ext in url.lower():  # Check if the URL contains the file extension
            file_type = ext.upper()  # Extract file type
            file_content = None
            if file_type == '.pdf':
                file_content = extract_pdf_content(canonical_url)
            elif file_type in ['.doc', '.docx']:
                file_content = extract_doc_content(canonical_url)
            elif file_type in ['.ppt', '.pptx']:
                file_content = extract_ppt_content(canonical_url)

            # Insert the file content into the page_data table
            if file_content:
                cur.execute("INSERT INTO page_data (page_id, data_type_code, data) VALUES (%s, %s, %s)",
                            (page_id, ext.upper(), file_content))
                conn.commit()

    images = extract_images(url, html)
    for image_url, image_content in images.items():
        save_image(page_id, image_url, image_content)


# Function to extract content from a PDF file
def extract_pdf_content(url):
    try:
        response = requests.get(url, timeout=5)
        pdf_file = PyPDF2.PdfFileReader(io.BytesIO(response.content))
        text = ""
        for page_num in range(pdf_file.numPages):
            text += pdf_file.getPage(page_num).extractText()
        return text
    except Exception as e:
        print(f"Error extracting PDF content from {url}: {e}")
        return None


# Function to extract content from a DOC or DOCX file
def extract_doc_content(url):
    try:
        response = requests.get(url, timeout=5)
        doc = docx.Document(io.BytesIO(response.content))
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text
        return text
    except Exception as e:
        print(f"Error extracting DOC/DOCX content from {url}: {e}")
        return None


# Function to extract content from a PPT or PPTX file
def extract_ppt_content(url):
    try:
        response = requests.get(url, timeout=5)
        ppt = pptx.Presentation(io.BytesIO(response.content))
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
        return text
    except Exception as e:
        print(f"Error extracting PPT/PPTX content from {url}: {e}")
        return None


def extract_images(site_url, html):
    soup = BeautifulSoup(html, 'html.parser')
    images = {}
    for img_tag in soup.find_all('img'):
        img_url = img_tag.get('src')
        if img_url:
            img_content = fetch_image(site_url, img_url)
            if img_content:
                images[img_url] = img_content
    return images


def fetch_image(site_url, image_url):
    try:
        if image_url.startswith('data:image'):
            # Extract base64 data from the URL
            base64_data = image_url.split(',')[1]
            # Decode base64 data
            image_content = base64.b64decode(base64_data)
            return image_content
        elif not urlparse(image_url).scheme:
            image_url = site_url + image_url  # Assuming HTTPS, you can adjust as needed
            response = requests.get(image_url, timeout=5)
            if response.status_code == 200:
                return response.content
    except Exception as e:
        print(f"Error fetching image from {image_url}: {e}")
    return None


def save_image(page_id, url, content):
    filename = url.split('/')[-1]  # Extract filename from URL
    content_type = 'image/jpeg'  # Adjust content type based on image format if needed
    try:
        cur.execute("INSERT INTO image (page_id, filename, content_type, data, accessed_time) "
                    "VALUES (%s, %s, %s, %s, %s)",
                    (page_id, filename, content_type, content, datetime.now()))
        conn.commit()
        # print(f"Image saved: {filename}")
    except Exception as e:
        print(f"Error saving image {filename}: {e}")


# Function to fetch and parse the robots.txt file
def parse_robots_txt(domain):
    robots_txt_url = urljoin(domain, "/robots.txt")
    try:
        response = requests.get(robots_txt_url, timeout=5)
        if response.status_code == 200:
            return response.text
    except Exception as e:
        print(f"Error fetching {robots_txt_url}: {e}")
    return ""


# Function to parse robots.txt rules
def parse_robots_rules(robots_txt_content):
    rules = {}
    current_user_agent = None
    for line in robots_txt_content.splitlines():
        if line.strip():  # Ignore empty lines
            if line.lower().startswith("user-agent"):
                current_user_agent = line.split(":")[1].strip()
            elif current_user_agent and ":" in line:  # Check if line contains a colon
                key, value = line.split(":", 1)
                rules.setdefault(current_user_agent, []).append((key.strip(), value.strip()))
    return rules


# Function to check if a URL is allowed to be crawled
def is_allowed_by_robots(url, robots_rules):
    parsed_url = urlparse(url)
    user_agent = "*"
    if "User-agent" in robots_rules:
        user_agent = robots_rules["User-agent"]
    elif "*" in robots_rules:
        user_agent = "*"
    if user_agent in robots_rules:
        for rule_key, rule_value in robots_rules[user_agent]:
            if rule_key.lower() == "allow" and rule_value != "/" and parsed_url.path.startswith(rule_value):
                return True
            elif rule_key.lower() == "disallow" and parsed_url.path.startswith(rule_value):
                return False
    return True


# Function to fetch and parse the sitemap.xml file
def parse_sitemap_xml(domain):
    sitemap_url = urljoin(domain, "/sitemap.xml")
    try:
        response = requests.get(sitemap_url, timeout=5)
        if response.status_code == 200:
            return response.text
    except Exception as e:
        print(f"Error fetching {sitemap_url}: {e}")
    return ""


def save_site_data(domain, robots_content, sitemap_content):
    try:
        # Check if the record already exists
        cur.execute("SELECT id FROM site WHERE domain = %s AND robots_content = %s AND sitemap_content = %s",
                    (domain, robots_content, sitemap_content))
        existing_record = cur.fetchone()

        if existing_record:
            print(f"Record for {domain} already exists in the database.")
            return existing_record[0]  # Return the existing record's ID

        # Insert the record as it doesn't exist
        cur.execute("INSERT INTO site (domain, robots_content, sitemap_content) VALUES (%s, %s, %s) "
                    "RETURNING id",
                    (domain, robots_content, sitemap_content))
        id = cur.fetchone()[0]
        conn.commit()
        return id
    except Exception as e:
        print(f"Error saving data for {domain}: {e}")


# Function to crawl websites
def crawl():
    global current_url, from_url
    while True:
        url = None
        # Ensure thread safety while accessing shared frontier
        cur.execute("DELETE FROM frontier "
                    "WHERE id = (SELECT id FROM frontier ORDER BY id LIMIT 1) "
                    "RETURNING link, from_link")
        row = cur.fetchone()
        try:
            url = row[0]
            current_url = row[0]
            from_url = row[1]
            domain = urlparse(url).netloc
            respect_crawl_delay(domain)
            robots = parse_robots_txt(url)
            rules = parse_robots_rules(robots)
            if is_allowed_by_robots(url, rules):
                sitemap = parse_sitemap_xml(url)
                site_id = save_site_data(domain, robots, sitemap)
                cur.execute("SELECT id FROM site WHERE domain = %s", (domain, ))
                crawl_url(url, site_id)
        except TypeError:
            print("Caught a TypeError: 'NoneType' object is not subscriptable")
        except Exception as e:
            print("Unknown error: " + str(e))


# Initialize database connection
conn = psycopg2.connect(database="postgres", user="postgres.guoimnempzxzidvwjnem", password="IepsCrawler12!!",
                        host="aws-0-eu-west-2.pooler.supabase.com", port="5432")

cur = conn.cursor()

# Global variable to store last access times for domains
global last_access_times
last_access_times = {}

global current_url
global from_url

# Number of threads to use
num_threads = 5  # Adjust as needed

# Queue to hold URLs to be processed
url_queue = Queue()


# Worker function for each thread
def worker():
    while True:
        url = url_queue.get()
        if url is None:
            break
        crawl_url(url)
        url_queue.task_done()


# Create and start threads
threads = []
for _ in range(num_threads):
    t = threading.Thread(target=worker)
    t.start()
    threads.append(t)

# Add URLs to the queue
# Example: url_queue.put('https://example.com')
# Add your URLs accordingly

# Wait for all threads to finish
url_queue.join()

# Stop worker threads
for _ in range(num_threads):
    url_queue.put(None)

# Join all threads
for t in threads:
    t.join()


if __name__ == "__main__":
    # n_threads = input("Please input number of threads: ")
    num_threads = int(1)
    crawl()
